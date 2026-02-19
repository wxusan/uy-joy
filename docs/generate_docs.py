#!/usr/bin/env python3
"""
Generate Technical Report (PDF) and Presentation (PPTX) for Uy-Joy Platform
"""

from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch, cm
from reportlab.lib.colors import HexColor
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

# Colors from design system
NAVY_900 = HexColor("#1E2A38")
GOLD_400 = HexColor("#C9A86A")
BACKGROUND = HexColor("#F7F8FA")
GREEN = HexColor("#4CAF50")
AMBER = HexColor("#F9A825")
RED = HexColor("#E53935")

def create_pdf_report():
    """Create detailed technical report PDF"""
    doc = SimpleDocTemplate(
        "Uy-Joy_Technical_Report.pdf",
        pagesize=A4,
        rightMargin=2*cm,
        leftMargin=2*cm,
        topMargin=2*cm,
        bottomMargin=2*cm
    )
    
    styles = getSampleStyleSheet()
    
    # Custom styles
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=NAVY_900,
        spaceAfter=30,
        alignment=TA_CENTER,
        fontName='Helvetica-Bold'
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=NAVY_900,
        spaceBefore=20,
        spaceAfter=10,
        fontName='Helvetica-Bold'
    )
    
    subheading_style = ParagraphStyle(
        'CustomSubheading',
        parent=styles['Heading3'],
        fontSize=12,
        textColor=NAVY_900,
        spaceBefore=15,
        spaceAfter=8,
        fontName='Helvetica-Bold'
    )
    
    body_style = ParagraphStyle(
        'CustomBody',
        parent=styles['Normal'],
        fontSize=10,
        textColor=HexColor("#374151"),
        spaceAfter=8,
        alignment=TA_JUSTIFY,
        leading=14
    )
    
    bullet_style = ParagraphStyle(
        'CustomBullet',
        parent=styles['Normal'],
        fontSize=10,
        textColor=HexColor("#374151"),
        leftIndent=20,
        spaceAfter=4,
        leading=14
    )
    
    story = []
    
    # Title Page
    story.append(Spacer(1, 2*inch))
    story.append(Paragraph("UY-JOY", title_style))
    story.append(Paragraph("Real Estate Visualization Platform", ParagraphStyle(
        'Subtitle',
        fontSize=14,
        textColor=GOLD_400,
        alignment=TA_CENTER,
        spaceAfter=30
    )))
    story.append(Spacer(1, 0.5*inch))
    story.append(Paragraph("Technical Documentation & Development Report", ParagraphStyle(
        'DocTitle',
        fontSize=18,
        textColor=NAVY_900,
        alignment=TA_CENTER,
        spaceAfter=50
    )))
    story.append(Spacer(1, 2*inch))
    story.append(Paragraph(f"Version 1.0 | {datetime.now().strftime('%B %d, %Y')}", ParagraphStyle(
        'Version',
        fontSize=10,
        textColor=HexColor("#6B7280"),
        alignment=TA_CENTER
    )))
    story.append(PageBreak())
    
    # Table of Contents
    story.append(Paragraph("Table of Contents", heading_style))
    toc_items = [
        "1. Executive Summary",
        "2. Technology Stack",
        "3. System Architecture",
        "4. Features Implemented",
        "5. Database Schema",
        "6. API Endpoints",
        "7. Design System",
        "8. Security & Performance",
        "9. Deployment Guide",
        "10. Future Roadmap"
    ]
    for item in toc_items:
        story.append(Paragraph(item, body_style))
    story.append(PageBreak())
    
    # 1. Executive Summary
    story.append(Paragraph("1. Executive Summary", heading_style))
    story.append(Paragraph(
        "Uy-Joy is a modern real estate visualization platform designed for property developers "
        "in Uzbekistan. The platform enables interactive floor plan exploration, apartment browsing, "
        "and lead generation through an intuitive user interface.",
        body_style
    ))
    story.append(Paragraph("Key Achievements:", subheading_style))
    achievements = [
        "• Full-stack Next.js 14 application with TypeScript",
        "• Multi-language support (Uzbek, English, Russian) with Uzbek as default",
        "• Interactive floor plan editor with drag-and-drop unit placement",
        "• Real-time apartment filtering and search",
        "• Integrated contact system (Telegram, Phone, Web forms)",
        "• Image optimization pipeline for fast loading",
        "• Responsive design for all devices"
    ]
    for a in achievements:
        story.append(Paragraph(a, bullet_style))
    story.append(PageBreak())
    
    # 2. Technology Stack
    story.append(Paragraph("2. Technology Stack", heading_style))
    
    tech_data = [
        ["Category", "Technology", "Version"],
        ["Frontend", "Next.js", "14.2.35"],
        ["Frontend", "React", "18.x"],
        ["Frontend", "TypeScript", "5.x"],
        ["Frontend", "Tailwind CSS", "3.4.x"],
        ["Backend", "Next.js API Routes", "14.x"],
        ["Database", "PostgreSQL", "15+"],
        ["ORM", "Prisma", "5.x"],
        ["i18n", "next-intl", "3.x"],
        ["Image Processing", "Sharp", "0.33.x"],
        ["Authentication", "NextAuth.js", "4.x"],
    ]
    
    tech_table = Table(tech_data, colWidths=[2*inch, 2.5*inch, 1.5*inch])
    tech_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), NAVY_900),
        ('TEXTCOLOR', (0, 0), (-1, 0), HexColor("#FFFFFF")),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 10),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), BACKGROUND),
        ('TEXTCOLOR', (0, 1), (-1, -1), HexColor("#374151")),
        ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
        ('FONTSIZE', (0, 1), (-1, -1), 9),
        ('GRID', (0, 0), (-1, -1), 1, HexColor("#E5E7EB")),
        ('TOPPADDING', (0, 1), (-1, -1), 8),
        ('BOTTOMPADDING', (0, 1), (-1, -1), 8),
    ]))
    story.append(tech_table)
    story.append(PageBreak())
    
    # 3. System Architecture
    story.append(Paragraph("3. System Architecture", heading_style))
    story.append(Paragraph("Application Structure:", subheading_style))
    
    arch_text = """
    The application follows a modern monolithic architecture using Next.js App Router:
    
    src/
    ├── app/                    # Next.js App Router pages
    │   ├── api/               # REST API endpoints
    │   ├── kvartiralar/       # Apartments listing page
    │   ├── kvartiralarni-korish/  # Explore page
    │   ├── portal/            # Admin management portal
    │   └── projects/          # Project detail pages
    ├── components/            # Reusable React components
    ├── lib/                   # Utilities and configurations
    └── messages/              # i18n translation files
    """
    story.append(Paragraph(arch_text.replace('\n', '<br/>'), body_style))
    
    story.append(Paragraph("Data Flow:", subheading_style))
    data_flow = [
        "• User requests page → Next.js Server Component fetches data",
        "• Prisma ORM queries PostgreSQL database",
        "• Data returned and rendered with React Server Components",
        "• Client components handle interactivity (filters, modals, forms)",
        "• Form submissions → API routes → Database updates"
    ]
    for item in data_flow:
        story.append(Paragraph(item, bullet_style))
    story.append(PageBreak())
    
    # 4. Features Implemented
    story.append(Paragraph("4. Features Implemented", heading_style))
    
    story.append(Paragraph("4.1 Public Features", subheading_style))
    public_features = [
        "• Homepage with animated statistics and featured apartments",
        "• Apartment listing page with advanced filters (rooms, area, price, status)",
        "• Interactive floor plan exploration (Building → Floor → Unit selection)",
        "• Unit detail modal with sketch image, specifications, and pricing",
        "• Contact form with lead capture (name + phone)",
        "• Floating contact sidebar (Telegram message, Phone call)",
        "• Multi-language support with automatic browser detection",
        "• Responsive design for mobile, tablet, and desktop"
    ]
    for f in public_features:
        story.append(Paragraph(f, bullet_style))
    
    story.append(Paragraph("4.2 Admin Features", subheading_style))
    admin_features = [
        "• Secure admin portal with authentication",
        "• Project management (create, edit, delete)",
        "• Building management with floor configurations",
        "• Interactive floor plan editor with unit placement",
        "• Unit management (rooms, area, price, status, sketch upload)",
        "• Reservation tracking (customer name, phone, notes)",
        "• Image gallery management per building",
        "• User management for admin accounts"
    ]
    for f in admin_features:
        story.append(Paragraph(f, bullet_style))
    story.append(PageBreak())
    
    # 5. Database Schema
    story.append(Paragraph("5. Database Schema", heading_style))
    story.append(Paragraph(
        "The database uses PostgreSQL with Prisma ORM. Key models include:",
        body_style
    ))
    
    schema_data = [
        ["Model", "Key Fields", "Relationships"],
        ["Project", "id, name, description, address", "Has many Buildings"],
        ["Building", "id, name, floors, projectId", "Belongs to Project, Has many Floors"],
        ["Floor", "id, number, planImage, buildingId", "Belongs to Building, Has many Units"],
        ["Unit", "id, unitNumber, rooms, area, status, price", "Belongs to Floor"],
        ["Lead", "id, name, phone, unitId, createdAt", "Optional Unit reference"],
        ["User", "id, email, password, role", "Authentication"],
    ]
    
    schema_table = Table(schema_data, colWidths=[1.2*inch, 2.3*inch, 2.5*inch])
    schema_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), NAVY_900),
        ('TEXTCOLOR', (0, 0), (-1, 0), HexColor("#FFFFFF")),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('GRID', (0, 0), (-1, -1), 1, HexColor("#E5E7EB")),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ('BACKGROUND', (0, 1), (-1, -1), BACKGROUND),
    ]))
    story.append(schema_table)
    
    story.append(Paragraph("Unit Status Values:", subheading_style))
    status_data = [
        ["Status", "Color", "Description"],
        ["available", "Green (#4CAF50)", "Ready for sale"],
        ["reserved", "Amber (#F9A825)", "Customer interested, pending payment"],
        ["sold", "Red (#E53935)", "Transaction completed"],
    ]
    status_table = Table(status_data, colWidths=[1.5*inch, 2*inch, 2.5*inch])
    status_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), NAVY_900),
        ('TEXTCOLOR', (0, 0), (-1, 0), HexColor("#FFFFFF")),
        ('GRID', (0, 0), (-1, -1), 1, HexColor("#E5E7EB")),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
    ]))
    story.append(status_table)
    story.append(PageBreak())
    
    # 6. API Endpoints
    story.append(Paragraph("6. API Endpoints", heading_style))
    
    api_data = [
        ["Endpoint", "Method", "Description"],
        ["/api/projects", "GET, POST", "List/Create projects"],
        ["/api/projects/[id]", "GET, PUT, DELETE", "Project CRUD"],
        ["/api/buildings", "GET, POST", "List/Create buildings"],
        ["/api/buildings/[id]", "GET, PUT, DELETE", "Building CRUD"],
        ["/api/floors", "GET, POST", "List/Create floors"],
        ["/api/floors/[id]", "GET, PUT, DELETE", "Floor CRUD"],
        ["/api/units", "GET, POST", "List/Create units"],
        ["/api/units/[id]", "GET, PUT, DELETE", "Unit CRUD"],
        ["/api/leads", "POST", "Submit contact form"],
        ["/api/upload", "POST", "Image upload with optimization"],
        ["/api/auth/[...nextauth]", "GET, POST", "Authentication"],
    ]
    
    api_table = Table(api_data, colWidths=[2.2*inch, 1.3*inch, 2.5*inch])
    api_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), NAVY_900),
        ('TEXTCOLOR', (0, 0), (-1, 0), HexColor("#FFFFFF")),
        ('GRID', (0, 0), (-1, -1), 1, HexColor("#E5E7EB")),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
        ('BACKGROUND', (0, 1), (-1, -1), BACKGROUND),
    ]))
    story.append(api_table)
    story.append(PageBreak())
    
    # 7. Design System
    story.append(Paragraph("7. Design System", heading_style))
    
    story.append(Paragraph("7.1 Color Palette", subheading_style))
    color_data = [
        ["Name", "Hex Code", "Usage"],
        ["Navy 900 (Primary)", "#1E2A38", "Headers, buttons, text"],
        ["Gold 400 (Accent)", "#C9A86A", "CTAs, highlights, phone button"],
        ["Background", "#F7F8FA", "Page backgrounds, cards"],
        ["Available", "#4CAF50", "Available unit status"],
        ["Reserved", "#F9A825", "Reserved unit status"],
        ["Sold", "#E53935", "Sold unit status"],
    ]
    color_table = Table(color_data, colWidths=[1.8*inch, 1.5*inch, 2.7*inch])
    color_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), NAVY_900),
        ('TEXTCOLOR', (0, 0), (-1, 0), HexColor("#FFFFFF")),
        ('GRID', (0, 0), (-1, -1), 1, HexColor("#E5E7EB")),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('TOPPADDING', (0, 0), (-1, -1), 6),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
    ]))
    story.append(color_table)
    
    story.append(Paragraph("7.2 Typography", subheading_style))
    typo_data = [
        ["Element", "Font", "Weight"],
        ["Headings", "Poppins", "600-700 (Semibold/Bold)"],
        ["Body Text", "Inter", "400-500 (Regular/Medium)"],
    ]
    typo_table = Table(typo_data, colWidths=[2*inch, 2*inch, 2*inch])
    typo_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), NAVY_900),
        ('TEXTCOLOR', (0, 0), (-1, 0), HexColor("#FFFFFF")),
        ('GRID', (0, 0), (-1, -1), 1, HexColor("#E5E7EB")),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
    ]))
    story.append(typo_table)
    
    story.append(Paragraph("7.3 Components", subheading_style))
    comp_items = [
        "• Border Radius: 10px (rounded-btn class)",
        "• Shadow: 0 4px 6px rgba(0,0,0,0.1) (shadow-card)",
        "• Buttons: Navy background, white text, gold for phone CTA",
        "• Forms: Compact, minimal fields (name + phone only)",
        "• Cards: White background, subtle shadow, hover lift effect"
    ]
    for c in comp_items:
        story.append(Paragraph(c, bullet_style))
    story.append(PageBreak())
    
    # 8. Security & Performance
    story.append(Paragraph("8. Security & Performance", heading_style))
    
    story.append(Paragraph("8.1 Security Measures", subheading_style))
    security = [
        "• Protected admin routes with NextAuth.js authentication",
        "• Hidden admin URL path (/portal/management-x7k9)",
        "• CSRF protection via Next.js built-in mechanisms",
        "• Input validation on all API endpoints",
        "• Environment variables for sensitive configuration"
    ]
    for s in security:
        story.append(Paragraph(s, bullet_style))
    
    story.append(Paragraph("8.2 Performance Optimizations", subheading_style))
    perf = [
        "• Image optimization with Sharp (resize, compress on upload)",
        "• Next.js Image component for automatic optimization",
        "• Server Components for reduced JavaScript bundle",
        "• Database query optimization with Prisma includes",
        "• Static generation where possible"
    ]
    for p in perf:
        story.append(Paragraph(p, bullet_style))
    
    story.append(Paragraph("Image Optimization Settings:", subheading_style))
    img_data = [
        ["Type", "Max Dimensions", "Quality"],
        ["Project/Building", "1200 x 800px", "JPEG 80%"],
        ["Floor Plan", "2000 x 1500px", "PNG Level 9"],
        ["Unit Sketch", "1200 x 1200px", "WebP 80%"],
    ]
    img_table = Table(img_data, colWidths=[2*inch, 2*inch, 2*inch])
    img_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), NAVY_900),
        ('TEXTCOLOR', (0, 0), (-1, 0), HexColor("#FFFFFF")),
        ('GRID', (0, 0), (-1, -1), 1, HexColor("#E5E7EB")),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
    ]))
    story.append(img_table)
    story.append(PageBreak())
    
    # 9. Deployment Guide
    story.append(Paragraph("9. Deployment Guide", heading_style))
    
    story.append(Paragraph("9.1 Prerequisites", subheading_style))
    prereqs = [
        "• Node.js 18+ installed",
        "• PostgreSQL 15+ database",
        "• npm or yarn package manager"
    ]
    for p in prereqs:
        story.append(Paragraph(p, bullet_style))
    
    story.append(Paragraph("9.2 Environment Variables", subheading_style))
    env_data = [
        ["Variable", "Description"],
        ["DATABASE_URL", "PostgreSQL connection string"],
        ["NEXTAUTH_SECRET", "Random secret for session encryption"],
        ["NEXTAUTH_URL", "Production URL of the application"],
    ]
    env_table = Table(env_data, colWidths=[2.5*inch, 3.5*inch])
    env_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), NAVY_900),
        ('TEXTCOLOR', (0, 0), (-1, 0), HexColor("#FFFFFF")),
        ('GRID', (0, 0), (-1, -1), 1, HexColor("#E5E7EB")),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
    ]))
    story.append(env_table)
    
    story.append(Paragraph("9.3 Deployment Steps", subheading_style))
    deploy_steps = [
        "1. Clone repository: git clone [repo-url]",
        "2. Install dependencies: npm install",
        "3. Configure environment: cp .env.example .env",
        "4. Run database migrations: npx prisma migrate deploy",
        "5. Build application: npm run build",
        "6. Start production server: npm start",
        "7. (Optional) Use PM2 for process management"
    ]
    for d in deploy_steps:
        story.append(Paragraph(d, bullet_style))
    story.append(PageBreak())
    
    # 10. Future Roadmap
    story.append(Paragraph("10. Future Roadmap", heading_style))
    
    story.append(Paragraph("10.1 Planned Features", subheading_style))
    planned = [
        "• AI Chatbot integration for customer support",
        "• Telegram Bot for notifications and inquiries",
        "• Comparison tool (compare multiple apartments)",
        "• Favorites/Wishlist functionality",
        "• PDF export for apartment details",
        "• Virtual tour / 3D visualization links",
        "• Analytics dashboard for admins",
        "• Bulk operations for unit management"
    ]
    for p in planned:
        story.append(Paragraph(p, bullet_style))
    
    story.append(Paragraph("10.2 Technical Improvements", subheading_style))
    tech_improve = [
        "• Implement Redis caching for frequently accessed data",
        "• Add comprehensive test suite (Jest, Playwright)",
        "• Set up CI/CD pipeline with GitHub Actions",
        "• Implement rate limiting on API endpoints",
        "• Add error tracking with Sentry"
    ]
    for t in tech_improve:
        story.append(Paragraph(t, bullet_style))
    
    story.append(Spacer(1, inch))
    story.append(Paragraph("— End of Technical Report —", ParagraphStyle(
        'End',
        fontSize=10,
        textColor=HexColor("#6B7280"),
        alignment=TA_CENTER
    )))
    
    doc.build(story)
    print("✅ PDF Report generated: Uy-Joy_Technical_Report.pdf")


def create_pptx_presentation():
    """Create presentation in landscape/album mode"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    
    # Helper to add slide with title
    def add_title_slide(title, subtitle=""):
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Background
        background = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(0x1E, 0x2A, 0x38)
        background.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.CENTER
        
        if subtitle:
            p = tf.add_paragraph()
            p.text = subtitle
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(0xC9, 0xA8, 0x6A)
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, bullet_points):
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Title bar
        title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.3))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = RGBColor(0x1E, 0x2A, 0x38)
        title_bar.line.fill.background()
        
        # Title text
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12), Inches(5.3))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(22)
            p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
            p.space_after = Pt(12)
        
        return slide
    
    def add_two_column_slide(title, left_items, right_items, left_title="", right_title=""):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        
        # Title bar
        title_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.3))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = RGBColor(0x1E, 0x2A, 0x38)
        title_bar.line.fill.background()
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.7))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        
        # Left column
        if left_title:
            lt_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(0.5))
            tf = lt_box.text_frame
            p = tf.paragraphs[0]
            p.text = left_title
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xC9, 0xA8, 0x6A)
        
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(6), Inches(5))
        tf = left_box.text_frame
        for i, item in enumerate(left_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
            p.space_after = Pt(8)
        
        # Right column
        if right_title:
            rt_box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(6), Inches(0.5))
            tf = rt_box.text_frame
            p = tf.paragraphs[0]
            p.text = right_title
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0xC9, 0xA8, 0x6A)
        
        right_box = slide.shapes.add_textbox(Inches(7), Inches(2.1), Inches(6), Inches(5))
        tf = right_box.text_frame
        for i, item in enumerate(right_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0x37, 0x41, 0x51)
            p.space_after = Pt(8)
        
        return slide
    
    # Slide 1: Title
    add_title_slide("UY-JOY", "Kvartiralarni onlayn sotish platformasi")
    
    # Slide 2: Problem & Solution
    add_content_slide("Muammo va yechim", [
        "Mijozlar kvartira tanlash uchun ofisga kelishi kerak",
        "Qog'oz kataloglar eskirgan va noqulay",
        "Savdo bo'limi har bir mijozga alohida vaqt sarflaydi",
        "Yechim: Onlayn platforma orqali 24/7 kvartira ko'rish",
        "Mijozlar o'zlari tanlab, keyin bog'lanadi"
    ])
    
    # Slide 3: Benefits for Company
    add_two_column_slide(
        "Kompaniya uchun foydalari",
        ["Savdo jarayoni tezlashadi", "Kamroq xodim vaqti sarflanadi", "Ko'proq mijozlarga xizmat", "Professional imidj", "Raqobatchilardan ajralib turish"],
        ["Barcha ma'lumotlar bir joyda", "Mijoz bazasi avtomatik to'planadi", "Har qanday qurilmada ishlaydi", "O'zbek, Rus, Ingliz tillarida", "Telegram orqali tezkor aloqa"],
        "Biznes uchun",
        "Qulayliklar"
    )
    
    # Slide 4: Benefits for Customers
    add_content_slide("Mijozlar uchun foydalari", [
        "Uydan chiqmasdan kvartiralarni ko'rish",
        "Narxlar, maydon, xonalar soni - barchasi ochiq",
        "Qaysi kvartira bo'sh, qaysi sotilgan - aniq ko'rinadi",
        "Bir tugma bilan Telegram yoki telefon orqali bog'lanish",
        "Telefondan ham qulay ishlaydi"
    ])
    
    # Slide 5: How it works
    add_content_slide("Qanday ishlaydi?", [
        "1. Mijoz saytga kiradi",
        "2. Binoni va qavatni tanlaydi",
        "3. Kvartira ustiga bosib tafsilotlarni ko'radi",
        "4. Yoqsa - Telegram yoki telefon orqali bog'lanadi",
        "5. Siz mijoz ma'lumotlarini admin panelda ko'rasiz"
    ])
    
    # Slide 6: Admin Panel
    add_content_slide("Boshqaruv paneli", [
        "Kvartiralar holatini o'zgartirish (Mavjud/Band/Sotilgan)",
        "Narxlarni yangilash",
        "Yangi bino va qavatlar qo'shish",
        "Qavat rejasi rasmlarini yuklash",
        "Mijoz so'rovlarini ko'rish",
        "Foydalanish juda oson - maxsus bilim talab qilinmaydi"
    ])
    
    # Slide 7: Contact Features
    add_content_slide("Aloqa imkoniyatlari", [
        "Har bir sahifada Telegram va Telefon tugmalari",
        "Mijoz bir bosish bilan bog'lanadi",
        "Ariza formasi - faqat ism va telefon (oddiy)",
        "Barcha so'rovlar bazada saqlanadi",
        "Kelajakda: AI chatbot qo'shiladi"
    ])
    
    # Slide 8: What's Next
    add_two_column_slide(
        "Kelajakda qo'shiladigan imkoniyatlar",
        ["AI yordamchi (chatbot)", "Telegram bot orqali xabarlar", "Kvartiralarni solishtirish", "Sevimlilar ro'yxati"],
        ["PDF formatda yuklab olish", "Sotuvlar statistikasi", "3D ko'rinish / Virtual tur", "Telegram orqali bildirishnomalar"],
        "Tez orada",
        "Keyingi bosqich"
    )
    
    # Slide 9: Summary
    add_title_slide("Hamkorlikka tayyormiz!", "Savollaringiz bo'lsa - bemalol so'rang")
    
    prs.save("Uy-Joy_Presentation.pptx")
    print("✅ PPTX Presentation generated: Uy-Joy_Presentation.pptx")


if __name__ == "__main__":
    create_pdf_report()
    create_pptx_presentation()
